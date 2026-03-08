#!/usr/bin/env python3
"""
将Markdown格式的PPT转换为PPTX文件
需要安装：pip install python-pptx
"""

import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def markdown_to_pptx(md_file, pptx_file):
    """将Markdown转换为PPTX"""
    
    # 读取Markdown文件
    with open(md_file, 'r', encoding='utf-8') as f:
        content = f.read()
    
    # 分割幻灯片
    slides = re.split(r'^## 幻灯片 \d+: ', content, flags=re.MULTILINE)
    
    # 创建PPT
    prs = Presentation()
    
    # 设置主题（使用空白布局）
    blank_slide_layout = prs.slide_layouts[6]  # 空白布局
    
    for i, slide_content in enumerate(slides):
        if i == 0:
            continue  # 跳过第一个空的部分
        
        # 提取标题和内容
        lines = slide_content.strip().split('\n')
        title = lines[0] if lines else f"幻灯片 {i}"
        content_lines = lines[1:] if len(lines) > 1 else []
        
        # 创建幻灯片
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # 添加标题
        if title:
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), 
                                                 Inches(9), Inches(1))
            title_frame = title_box.text_frame
            title_frame.text = title
            title_frame.paragraphs[0].font.size = Pt(32)
            title_frame.paragraphs[0].font.bold = True
            title_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
        
        # 添加内容
        if content_lines:
            content_text = '\n'.join(content_lines)
            
            # 清理Markdown格式
            content_text = re.sub(r'^### ', '', content_text, flags=re.MULTILINE)
            content_text = re.sub(r'^#### ', '', content_text, flags=re.MULTILINE)
            content_text = re.sub(r'\*\*(.*?)\*\*', r'\1', content_text)  # 移除粗体
            content_text = re.sub(r'\*(.*?)\*', r'\1', content_text)  # 移除斜体
            content_text = re.sub(r'`(.*?)`', r'\1', content_text)  # 移除代码标记
            
            # 处理表格
            content_text = re.sub(r'^\|.*\|$', '', content_text, flags=re.MULTILINE)
            
            content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), 
                                                   Inches(9), Inches(5))
            content_frame = content_box.text_frame
            content_frame.text = content_text
            
            # 设置字体
            for paragraph in content_frame.paragraphs:
                paragraph.font.size = Pt(18)
                paragraph.alignment = PP_ALIGN.LEFT
    
    # 保存PPTX
    prs.save(pptx_file)
    print(f"PPTX文件已保存: {pptx_file}")

def create_simple_pptx(md_file, pptx_file):
    """创建简化版PPTX（不依赖python-pptx）"""
    
    with open(md_file, 'r', encoding='utf-8') as f:
        content = f.read()
    
    # 生成HTML格式的PPT
    slides = re.split(r'^## 幻灯片 \d+: ', content, flags=re.MULTILINE)
    
    html_content = """
    <!DOCTYPE html>
    <html>
    <head>
        <meta charset="UTF-8">
        <title>芯片方案对比PPT</title>
        <style>
            body { font-family: Arial, sans-serif; margin: 0; padding: 0; }
            .slide { 
                width: 1024px; 
                height: 768px; 
                padding: 40px; 
                box-sizing: border-box;
                page-break-after: always;
                border-bottom: 2px solid #ccc;
            }
            .slide-title { 
                font-size: 32px; 
                font-weight: bold; 
                margin-bottom: 30px;
                color: #2c3e50;
            }
            .slide-content { 
                font-size: 18px; 
                line-height: 1.6;
                color: #34495e;
            }
            table {
                border-collapse: collapse;
                width: 100%;
                margin: 20px 0;
            }
            th, td {
                border: 1px solid #ddd;
                padding: 8px;
                text-align: left;
            }
            th {
                background-color: #f2f2f2;
                font-weight: bold;
            }
            .highlight { color: #e74c3c; font-weight: bold; }
            .check { color: #27ae60; }
            .cross { color: #c0392b; }
            @media print {
                .slide { 
                    width: 100%; 
                    height: 100vh; 
                }
            }
        </style>
    </head>
    <body>
    """
    
    for i, slide_content in enumerate(slides):
        if i == 0:
            continue
        
        lines = slide_content.strip().split('\n')
        title = lines[0] if lines else f"幻灯片 {i}"
        content_lines = lines[1:] if len(lines) > 1 else []
        
        html_content += f'<div class="slide">\n'
        html_content += f'<div class="slide-title">{title}</div>\n'
        html_content += '<div class="slide-content">\n'
        
        for line in content_lines:
            # 处理特殊格式
            line = line.replace('✅', '<span class="check">✅</span>')
            line = line.replace('❌', '<span class="cross">❌</span>')
            line = line.replace('⚠️', '<span class="highlight">⚠️</span>')
            line = line.replace('🥇', '<span class="highlight">🥇</span>')
            line = line.replace('🥈', '<span class="highlight">🥈</span>')
            
            # 处理表格行
            if line.startswith('|') and line.endswith('|'):
                if '---' in line:
                    continue  # 跳过表格分隔线
                cells = line.split('|')[1:-1]
                if i == 1:  # 第一行可能是表头
                    html_content += '<table>\n<thead>\n<tr>\n'
                    for cell in cells:
                        html_content += f'<th>{cell.strip()}</th>\n'
                    html_content += '</tr>\n</thead>\n<tbody>\n'
                else:
                    html_content += '<tr>\n'
                    for cell in cells:
                        html_content += f'<td>{cell.strip()}</td>\n'
                    html_content += '</tr>\n'
            elif line.strip():
                html_content += f'<p>{line}</p>\n'
        
        html_content += '</div>\n</div>\n'
    
    html_content += """
    </body>
    </html>
    """
    
    with open(pptx_file.replace('.pptx', '.html'), 'w', encoding='utf-8') as f:
        f.write(html_content)
    
    print(f"HTML格式PPT已保存: {pptx_file.replace('.pptx', '.html')}")
    print("可以在浏览器中打开并打印为PDF")

if __name__ == "__main__":
    md_file = "芯片对比PPT.md"
    pptx_file = "老年人智能助听耳塞-芯片方案对比.pptx"
    
    try:
        # 尝试使用python-pptx
        from pptx import Presentation
        markdown_to_pptx(md_file, pptx_file)
    except ImportError:
        print("python-pptx未安装，生成HTML格式PPT")
        create_simple_pptx(md_file, pptx_file)
    except Exception as e:
        print(f"生成PPTX时出错: {e}")
        print("生成HTML格式PPT作为备选")
        create_simple_pptx(md_file, pptx_file)